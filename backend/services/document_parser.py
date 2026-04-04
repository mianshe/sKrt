import base64
import io
import json
import os
import re
import subprocess
import sys
import time
from dataclasses import dataclass
from typing import Any, Dict, Iterator, List, Optional, Tuple

import httpx

from PyPDF2 import PdfReader
from docx import Document
from backend.services.glm_ocr import glm_ocr_enabled, run_glm_layout_parsing

try:
    import pdfplumber  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    pdfplumber = None  # type: ignore

try:
    from pdf2image import convert_from_path  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    convert_from_path = None  # type: ignore

try:
    import pytesseract  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    pytesseract = None  # type: ignore

try:
    import numpy as np  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    np = None  # type: ignore

# Before paddle/paddleocr import: reduce OneDNN/PIR issues on CPU (PaddlePaddle 3.3+).
# Full mitigation for PaddleOCR 3.x is enable_mkldnn=False in PaddleOCR(); see _get_paddle_runtime.
os.environ.setdefault("FLAGS_use_mkldnn", "0")
os.environ.setdefault("PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK", "1")

try:
    from paddleocr import PaddleOCR  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    PaddleOCR = None  # type: ignore

try:
    import pypdfium2 as pdfium  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    pdfium = None  # type: ignore

# Windows: child exit 3221225477 (unsigned) or -1073741819 (signed) == STATUS_ACCESS_VIOLATION (0xC0000005).
_WIN_STATUS_ACCESS_VIOLATION_U = 3221225477
_WIN_STATUS_ACCESS_VIOLATION_S = -1073741819


def _describe_subprocess_exit(code: Optional[int]) -> str:
    if code is None:
        return "unknown-exit"
    if code in {_WIN_STATUS_ACCESS_VIOLATION_U, _WIN_STATUS_ACCESS_VIOLATION_S}:
        return (
            "paddle_native_crash(ACCESS_VIOLATION_0xC0000005); "
            "install Tesseract and set PATH or TESSERACT_CMD, or PDF_OCR_ENGINE=tesseract; "
            "try lower PDF_OCR_DPI"
        )
    return f"exit-{code}"


@dataclass
class ParsedDocument:
    text: str
    metadata: Dict[str, Any]


class DocumentParser:
    def __init__(self) -> None:
        self._paddle_ocr: Optional[Any] = None
        self._paddle_init_error: Optional[str] = None
        self._baidu_token: Optional[str] = None
        self._baidu_token_deadline: float = 0.0
        self._ocr_runtime_metadata: Dict[str, Any] = {}
        # OCR 逐页缓存（由 upload_ingestion_service 注入 task_id；读写走 knowledge_store）
        self._ocr_cache_task_id: int = 0
        # 与供应商计费对齐：百度/HTTP OCR 成功 API 调用次数（入库后扣减余额）
        self._ocr_billable_api_calls: int = 0

    def set_ocr_cache(self, task_id: int) -> None:
        """注入 OCR 页级缓存任务 ID，由 UploadIngestionService 调用。"""
        self._ocr_cache_task_id = task_id

    def _ocr_cache_get(self, page_num: int) -> Optional[str]:
        if not self._ocr_cache_task_id:
            return None
        try:
            from backend.services import knowledge_store

            conn = knowledge_store.connect()
            try:
                row = conn.execute(
                    "SELECT ocr_text FROM ocr_page_cache WHERE task_id=? AND page_num=?",
                    (self._ocr_cache_task_id, page_num),
                ).fetchone()
                if not row:
                    return None
                return str(row["ocr_text"])
            finally:
                conn.close()
        except Exception:
            return None

    def _ocr_cache_set(self, page_num: int, text: str, engine: str = "") -> None:
        if not self._ocr_cache_task_id:
            return
        try:
            from backend.services import knowledge_store

            conn = knowledge_store.connect()
            try:
                if knowledge_store.use_postgres():
                    conn.execute(
                        """
                        INSERT INTO ocr_page_cache(task_id, page_num, ocr_text, engine)
                        VALUES (?, ?, ?, ?)
                        ON CONFLICT (task_id, page_num) DO UPDATE SET
                          ocr_text = EXCLUDED.ocr_text,
                          engine = EXCLUDED.engine
                        """,
                        (self._ocr_cache_task_id, page_num, text, engine),
                    )
                else:
                    conn.execute(
                        "INSERT OR REPLACE INTO ocr_page_cache(task_id,page_num,ocr_text,engine) VALUES(?,?,?,?)",
                        (self._ocr_cache_task_id, page_num, text, engine),
                    )
                conn.commit()
            finally:
                conn.close()
        except Exception:
            pass

    def parse(self, file_path: str, document_type: str, *, ocr_engine_override: Optional[str] = None) -> ParsedDocument:
        self._ocr_billable_api_calls = 0
        self._ocr_runtime_metadata = {}
        ext = os.path.splitext(file_path)[1].lower()
        pdf_page_count: Optional[int] = None
        if ext == ".pdf":
            text = self._parse_pdf(file_path, ocr_engine_override=ocr_engine_override)
            pdf_page_count = self._pdf_page_count(file_path)
        elif ext == ".docx":
            text = self._parse_docx(file_path)
        elif ext == ".pptx":
            text = self._parse_pptx(file_path)
        elif ext in {".txt", ".md", ".markdown", ".py", ".json"}:
            text = self._parse_text(file_path)
        elif ext in {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}:
            text = self._parse_image_file(file_path)
        else:
            raise ValueError(f"不支持的文件类型: {ext}")

        text = text.strip()
        metadata = self._extract_metadata(text, document_type, os.path.basename(file_path))
        if pdf_page_count is not None and pdf_page_count > 0:
            metadata["pdf_page_count"] = int(pdf_page_count)
        if ext == ".pdf":
            metadata["ocr_billable_api_calls"] = int(max(0, self._ocr_billable_api_calls))
            metadata["toc"] = self._extract_pdf_toc(file_path)
        metadata.update(self._ocr_runtime_metadata)

        # 自动检测文档形态、图片、编码质量
        metadata["document_form"] = self._infer_document_form(
            text, os.path.basename(file_path), page_count=pdf_page_count or 0
        )
        has_images, image_count = self._detect_has_images(file_path, ext)
        metadata["has_images"] = has_images
        metadata["image_count_estimate"] = image_count
        enc = self._check_encoding_quality(text)
        metadata["encoding_ok"] = enc["ok"]
        metadata["encoding_issues"] = enc["issues"]

        return ParsedDocument(text=text, metadata=metadata)

    def _parse_pdf(self, file_path: str, *, ocr_engine_override: Optional[str] = None) -> str:
        override = (ocr_engine_override or "").strip().lower()
        if override in {"glm-ocr", "glm_ocr", "complex_layout"}:
            return self._parse_pdf_with_glm_ocr(file_path)

        candidates: List[str] = []
        errors: List[str] = []

        if pdfplumber is not None:
            try:
                with pdfplumber.open(file_path) as pdf:
                    parts: List[str] = []
                    for i, p in enumerate(pdf.pages, 1):
                        pt = (p.extract_text() or "").strip()
                        if pt:
                            parts.append("[[PAGE:" + str(i) + "]]")
                            parts.append(pt)
                    txt = self._normalize_pdf_text(chr(10).join(parts))
                    if txt:
                        candidates.append(txt)
                    else:
                        errors.append("pdfplumber=empty-text")
            except Exception as exc:
                errors.append(f"pdfplumber={exc}")

        try:
            reader = PdfReader(file_path)
            parts = []
            for i, page in enumerate(reader.pages, 1):
                pt = (page.extract_text() or "").strip()
                if pt:
                    parts.append("[[PAGE:" + str(i) + "]]")
                    parts.append(pt)
            txt = self._normalize_pdf_text(chr(10).join(parts))
            if txt:
                candidates.append(txt)
            else:
                errors.append("pypdf2=empty-text")
        except Exception as exc:
            errors.append(f"pypdf2={exc}")

        if candidates:
            candidates.sort(key=lambda x: self._text_quality_score(x), reverse=True)
            best = candidates[0]
            if self._prefer_direct_pdf_text(best):
                return best

        ocr_text, ocr_error = self._ocr_pdf(file_path, ocr_engine_override=ocr_engine_override)
        if ocr_text:
            return ocr_text
        if ocr_error:
            errors.append(f"ocr={ocr_error}")

        _err_msg = "; ".join(errors) if errors else "unknown"
        raise ValueError(f"PDF解析失败: {_err_msg}")

    def _parse_pdf_with_glm_ocr(self, file_path: str) -> str:
        if not glm_ocr_enabled():
            raise ValueError("GLM-OCR 未配置 ZHIPU_API_KEY")
        result = run_glm_layout_parsing(file_path)
        self._ocr_runtime_metadata.update(
            {
                "ocr_used": True,
                "ocr_engine": "glm-ocr",
                "ocr_provider": "zhipu",
                "ocr_billable_tokens": int(max(0, result.total_tokens)),
                "ocr_prompt_tokens": int(max(0, result.prompt_tokens)),
                "ocr_completion_tokens": int(max(0, result.completion_tokens)),
                "ocr_request_id": result.request_id,
                "ocr_complex_layout": True,
            }
        )
        return result.text

    def _ocr_pdf(self, file_path: str, *, ocr_engine_override: Optional[str] = None) -> Tuple[str, str]:
        ocr_engine = (ocr_engine_override or os.getenv("PDF_OCR_ENGINE", "auto")).strip().lower() or "auto"
        ocr_dpi = self._safe_int(os.getenv("PDF_OCR_DPI", "260"), 260, min_value=120, max_value=500)
        ocr_max_pages = self._safe_int(os.getenv("PDF_OCR_MAX_PAGES", "0"), 0, min_value=0, max_value=5000)

        # 先检查是否至少有一页可渲染（避免空迭代器静默失败）
        if self._pdf_page_count(file_path) <= 0:
            return "", "render-failed:page-count-unavailable"

        engines: List[str]
        if ocr_engine == "paddle":
            engines = ["paddle"]
        elif ocr_engine == "tesseract":
            engines = ["tesseract"]
        elif ocr_engine in {"local", "local-standard", "paddle-local"}:
            engines = ["paddle", "tesseract"]
        elif ocr_engine == "baidu":
            engines = ["baidu"]
        elif ocr_engine in {"runpod", "gpu", "remote", "ocr_api", "api"}:
            engines = ["remote"]
        else:
            # auto：已配置百度时默认「联网 OCR 优先」，失败再 Paddle/Tesseract。
            # PDF_OCR_REMOTE_FIRST=0 则恢复本地优先、百度仅作兜底（PDF_OCR_BAIDU_FALLBACK=1）。
            baidu_ok = self._baidu_ocr_configured()
            remote_first = baidu_ok and self._as_bool(os.getenv("PDF_OCR_REMOTE_FIRST", "1"), True)
            baidu_fallback = baidu_ok and self._as_bool(os.getenv("PDF_OCR_BAIDU_FALLBACK", "1"), True)
            if remote_first:
                engines = ["baidu", "paddle", "tesseract"]
            elif baidu_fallback:
                engines = ["paddle", "tesseract", "baidu"]
            else:
                engines = ["paddle", "tesseract"]

        engine_errors: List[str] = []
        for engine in engines:
            if engine == "paddle":
                txt, err = self._ocr_pdf_pages_paddle(file_path=file_path, dpi=ocr_dpi, max_pages=ocr_max_pages)
                if txt:
                    self._ocr_runtime_metadata.update(
                        {"ocr_used": True, "ocr_engine": "paddle", "ocr_provider": "local", "ocr_call_billing_units": 1}
                    )
            elif engine == "tesseract":
                txt, err = self._ocr_pdf_pages_tesseract(file_path=file_path, dpi=ocr_dpi, max_pages=ocr_max_pages)
                if txt:
                    self._ocr_runtime_metadata.update(
                        {"ocr_used": True, "ocr_engine": "tesseract", "ocr_provider": "local", "ocr_call_billing_units": 1}
                    )
            elif engine == "remote":
                txt, err, ncalls = self._ocr_pdf_pages_http_api(file_path=file_path, dpi=ocr_dpi, max_pages=ocr_max_pages)
                if txt:
                    self._ocr_billable_api_calls = ncalls
                    self._ocr_runtime_metadata.update({"ocr_used": True, "ocr_engine": "remote-http", "ocr_provider": "remote"})
                    return txt, ""
                engine_errors.append(f"{engine}:{err or 'empty'}")
                continue
            else:
                txt, err, ncalls = self._ocr_pdf_pages_baidu(file_path=file_path, dpi=ocr_dpi, max_pages=ocr_max_pages)
                if txt:
                    self._ocr_billable_api_calls = ncalls
                    self._ocr_runtime_metadata.update({"ocr_used": True, "ocr_engine": "baidu", "ocr_provider": "baidu"})
                    return txt, ""
                engine_errors.append(f"{engine}:{err or 'empty'}")
                continue
            if txt:
                return txt, ""
            engine_errors.append(f"{engine}:{err or 'empty'}")

        return "", "; ".join(engine_errors)

    def _ocr_pdf_pages_http_api(self, file_path: str, dpi: int, max_pages: int) -> Tuple[str, str, int]:
        """
        外部 HTTP OCR（例如远端 PaddleOCR 服务）。优先读 OCR_API_*，仍兼容 GPU_OCR_*。
          - OCR_API_BASE：根 URL，无尾斜杠（旧：GPU_OCR_ENDPOINT）
          - OCR_API_KEY：可选，请求头 X-API-Key（旧：GPU_OCR_API_KEY）
          - OCR_API_TIMEOUT_SEC（旧：GPU_OCR_TIMEOUT_SEC）
        POST {OCR_API_BASE}/ocr/pdf，multipart file，JSON { ok, text, error? }。
        返回 (text, error, successful_http_calls)。
        """
        endpoint = (os.getenv("OCR_API_BASE", "") or os.getenv("GPU_OCR_ENDPOINT", "") or "").strip().rstrip("/")
        if not endpoint:
            return "", "ocr-api-base-not-configured", 0
        api_key = (os.getenv("OCR_API_KEY", "") or os.getenv("GPU_OCR_API_KEY", "") or "").strip()
        timeout_raw = (os.getenv("OCR_API_TIMEOUT_SEC", "") or os.getenv("GPU_OCR_TIMEOUT_SEC", "") or "900").strip() or "900"
        timeout_sec = self._safe_int(timeout_raw, 900, min_value=30, max_value=7200)
        url = f"{endpoint}/ocr/pdf"
        headers: Dict[str, str] = {}
        if api_key:
            headers["X-API-Key"] = api_key
        try:
            with open(file_path, "rb") as fp:
                files = {"file": (os.path.basename(file_path) or "document.pdf", fp, "application/pdf")}
                data = {"dpi": str(int(dpi)), "max_pages": str(int(max_pages))}
                resp = httpx.post(url, headers=headers, files=files, data=data, timeout=float(timeout_sec))
            resp.raise_for_status()
            payload = resp.json()
            if isinstance(payload, dict) and payload.get("ok") and str(payload.get("text") or "").strip():
                return str(payload.get("text") or ""), "", 1
            err = ""
            if isinstance(payload, dict):
                err = str(payload.get("error") or payload.get("detail") or "")
            return "", err or "ocr-api-empty", 0
        except Exception as exc:
            return "", f"ocr-api-failed:{exc}", 0

    def _pdf_page_count(self, file_path: str) -> int:
        try:
            reader = PdfReader(file_path)
            return len(reader.pages)
        except Exception:
            pass
        if pdfium is not None:
            try:
                document = pdfium.PdfDocument(file_path)
                n = len(document)
                document.close()
                return int(n)
            except Exception:
                pass
        return 0

    def _extract_pdf_toc(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Best-effort PDF outline extraction. Failures should not block parsing.
        """
        try:
            reader = PdfReader(file_path)
        except Exception:
            return []

        outlines: Any = None
        for attr_name in ("outline", "outlines"):
            try:
                outlines = getattr(reader, attr_name)
                if outlines:
                    break
            except Exception:
                outlines = None

        if not outlines:
            for method_name in ("get_outlines", "getOutlines"):
                method = getattr(reader, method_name, None)
                if callable(method):
                    try:
                        outlines = method()
                        if outlines:
                            break
                    except Exception:
                        outlines = None

        if not outlines:
            return []

        items: List[Dict[str, Any]] = []
        self._flatten_pdf_toc(reader, outlines, items, level=1)
        return items

    def _flatten_pdf_toc(
        self,
        reader: PdfReader,
        outlines: Any,
        items: List[Dict[str, Any]],
        *,
        level: int,
    ) -> None:
        if isinstance(outlines, list):
            for entry in outlines:
                if isinstance(entry, list):
                    self._flatten_pdf_toc(reader, entry, items, level=level + 1)
                    continue
                item = self._pdf_toc_entry_to_dict(reader, entry, level=level)
                if item is not None:
                    items.append(item)
            return

        item = self._pdf_toc_entry_to_dict(reader, outlines, level=level)
        if item is not None:
            items.append(item)

    def _pdf_toc_entry_to_dict(self, reader: PdfReader, entry: Any, *, level: int) -> Optional[Dict[str, Any]]:
        try:
            title = getattr(entry, "title", None)
            if title is None and isinstance(entry, dict):
                title = entry.get("/Title") or entry.get("title")
            title_text = str(title or "").strip()
            if not title_text:
                return None

            payload: Dict[str, Any] = {"title": title_text, "level": int(level)}
            page_num = self._resolve_pdf_toc_page(reader, entry)
            if page_num is not None:
                payload["page"] = int(page_num)
            return payload
        except Exception:
            return None

    def _resolve_pdf_toc_page(self, reader: PdfReader, entry: Any) -> Optional[int]:
        for method_name in ("get_destination_page_number", "getDestinationPageNumber"):
            method = getattr(reader, method_name, None)
            if callable(method):
                try:
                    page_index = method(entry)
                    if isinstance(page_index, int) and page_index >= 0:
                        return page_index + 1
                except Exception:
                    pass

        direct_page = None
        for attr_name in ("page", "page_number", "/Page"):
            try:
                if isinstance(entry, dict):
                    direct_page = entry.get(attr_name)
                else:
                    direct_page = getattr(entry, attr_name, None)
                if direct_page is not None:
                    break
            except Exception:
                pass

        if isinstance(direct_page, int) and direct_page >= 0:
            return direct_page + 1

        return None

    def _iter_pdf_page_images(self, file_path: str, dpi: int, max_pages: int) -> Iterator[Any]:
        """
        逐页产出 PIL.Image，避免一次性加载整本扫描件导致内存暴涨或进程被系统终止。
        max_pages==0 表示处理到最后一页。
        """
        try:
            page_total = len(PdfReader(file_path).pages)
        except Exception:
            page_total = 0
        if page_total <= 0 and pdfium is not None:
            try:
                document = pdfium.PdfDocument(file_path)
                page_total = len(document)
                document.close()
            except Exception:
                page_total = 0
        if page_total <= 0:
            return
        limit = min(page_total, max_pages) if max_pages > 0 else page_total
        scale = max(1.0, float(dpi) / 72.0)
        poppler_path = os.getenv("PDF_OCR_POPPLER_PATH", "").strip() or None

        if pdfium is not None:
            try:
                document = pdfium.PdfDocument(file_path)
                for i in range(limit):
                    page = document[i]
                    rendered = page.render(scale=scale)
                    pil_img = rendered.to_pil()
                    page.close()
                    yield pil_img
                document.close()
                return
            except Exception:
                pass

        if convert_from_path is not None:
            for p in range(1, limit + 1):
                imgs = convert_from_path(
                    file_path,
                    dpi=dpi,
                    first_page=p,
                    last_page=p,
                    poppler_path=poppler_path,
                )
                if imgs:
                    yield imgs[0]
            return

    def _ocr_pdf_pages_paddle(self, file_path: str, dpi: int, max_pages: int) -> Tuple[str, str]:
        # PDF_OCR_PADDLE_SUBPROCESS: 1/on (default) = run Paddle in child process (backend.services.ocr_worker);
        #   0/off = in-process (debug). PDF_OCR_SUBPROCESS_TIMEOUT: seconds for child run (default 3600; 0 = unlimited).
        if self._as_bool(os.getenv("PDF_OCR_PADDLE_SUBPROCESS", "1"), True):
            return self._ocr_pdf_pages_paddle_subprocess(file_path, dpi, max_pages)
        return self._ocr_pdf_pages_paddle_inprocess(file_path, dpi, max_pages)

    def _ocr_pdf_pages_paddle_subprocess(self, file_path: str, dpi: int, max_pages: int) -> Tuple[str, str]:
        timeout_sec = self._safe_int(os.getenv("PDF_OCR_SUBPROCESS_TIMEOUT", "3600"), 3600, min_value=0, max_value=86400 * 7)
        timeout = None if timeout_sec <= 0 else float(timeout_sec)
        safe_dpi = self._safe_int(os.getenv("PDF_OCR_PADDLE_SAFE_DPI", "180"), 180, min_value=120, max_value=300)
        attempt_specs: List[Tuple[int, Dict[str, str], str]] = [
            (dpi, {}, "default"),
            (
                min(dpi, safe_dpi),
                {
                    "PADDLEOCR_USE_ANGLE_CLS": "0",
                    "OMP_NUM_THREADS": "1",
                    "MKL_NUM_THREADS": "1",
                    "OPENBLAS_NUM_THREADS": "1",
                    "NUMEXPR_NUM_THREADS": "1",
                    "PADDLEOCR_ENABLE_MKLDNN": "0",
                },
                "safe-retry",
            ),
        ]
        errors: List[str] = []
        for run_dpi, overrides, tag in attempt_specs:
            txt, err, rc, crash = self._run_paddle_worker_once(
                file_path=file_path,
                dpi=run_dpi,
                max_pages=max_pages,
                timeout=timeout,
                timeout_sec=timeout_sec,
                env_overrides=overrides,
            )
            if txt:
                return txt, ""
            errors.append(f"{tag}:{err}")
            # First attempt non-crash often indicates real OCR/runtime issue; avoid pointless retry noise.
            if tag == "default" and not crash:
                break
        return "", " | ".join(errors)

    def _run_paddle_worker_once(
        self,
        file_path: str,
        dpi: int,
        max_pages: int,
        timeout: Optional[float],
        timeout_sec: int,
        env_overrides: Dict[str, str],
    ) -> Tuple[str, str, Optional[int], bool]:
        env = os.environ.copy()
        env.update(env_overrides)
        cmd = [
            sys.executable,
            "-m",
            "backend.services.ocr_worker",
            "paddle",
            file_path,
            str(dpi),
            str(max_pages),
        ]
        try:
            proc = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=timeout,
                env=env,
            )
        except subprocess.TimeoutExpired:
            return "", f"subprocess-timeout:{timeout_sec}s", None, False
        rc = proc.returncode
        stderr_tail = (proc.stderr or "").strip()[-4000:]
        stdout = (proc.stdout or "").strip()
        last_line = ""
        if stdout:
            lines = stdout.splitlines()
            last_line = lines[-1].strip() if lines else ""
        crash = rc in {_WIN_STATUS_ACCESS_VIOLATION_U, _WIN_STATUS_ACCESS_VIOLATION_S}
        try:
            data = json.loads(last_line) if last_line else {}
        except json.JSONDecodeError:
            base = _describe_subprocess_exit(rc)
            if not crash and (stderr_tail or stdout):
                base = f"{base}; tail={stderr_tail or stdout[:1200]}"
            return "", f"subprocess-bad-json:{base}", rc, crash

        if isinstance(data, dict) and data.get("ok") and isinstance(data.get("text"), str):
            return data["text"], "", rc, False

        err = ""
        if isinstance(data, dict):
            err = str(data.get("error") or "")
        if not err:
            err = _describe_subprocess_exit(rc)
        if stderr_tail and not crash:
            err = f"{err}; stderr={stderr_tail}"
        return "", err, rc, crash

    def _ocr_pdf_pages_paddle_inprocess(self, file_path: str, dpi: int, max_pages: int) -> Tuple[str, str]:
        if PaddleOCR is None:
            return "", "paddleocr-not-installed"
        if np is None:
            return "", "numpy-not-installed"
        runtime = self._get_paddle_runtime()
        if runtime is None:
            return "", f"runtime-init-failed:{self._paddle_init_error or 'unknown'}"
        ocr_pages: List[str] = []
        try:
            for img in self._iter_pdf_page_images(file_path, dpi, max_pages):
                try:
                    arr = np.array(img)
                    try:
                        result = runtime.ocr(arr, cls=True)
                    except TypeError as exc:
                        if "cls" in str(exc).lower() or "unexpected keyword" in str(exc).lower():
                            result = runtime.ocr(arr)
                        else:
                            raise
                    texts = self._extract_paddle_texts(result)
                    if texts:
                        ocr_pages.append("\n".join(texts))
                finally:
                    try:
                        img.close()
                    except Exception:
                        pass
        except Exception as exc:
            return "", f"ocr-failed:{exc}"
        merged = self._normalize_pdf_text("\n".join(ocr_pages))
        if not merged:
            return "", "empty-text"
        return merged, ""

    def _ocr_pdf_pages_tesseract(self, file_path: str, dpi: int, max_pages: int) -> Tuple[str, str]:
        if pytesseract is None:
            return "", "pytesseract-not-installed"
        tesseract_cmd = os.getenv("TESSERACT_CMD", "").strip()
        if tesseract_cmd:
            try:
                pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
            except Exception:
                pass
        ocr_lang = os.getenv("PDF_OCR_TESSERACT_LANG", os.getenv("PDF_OCR_LANG", "chi_sim+eng")).strip() or "chi_sim+eng"
        ocr_config = os.getenv("PDF_OCR_TESSERACT_CONFIG", "").strip()
        ocr_pages: List[str] = []
        try:
            for img in self._iter_pdf_page_images(file_path, dpi, max_pages):
                try:
                    page_text = (
                        pytesseract.image_to_string(img, lang=ocr_lang, config=ocr_config)
                        if ocr_config
                        else pytesseract.image_to_string(img, lang=ocr_lang)
                    )
                    if page_text and page_text.strip():
                        ocr_pages.append(page_text)
                finally:
                    try:
                        img.close()
                    except Exception:
                        pass
        except Exception as exc:
            return "", f"ocr-failed:{exc}"
        merged = self._normalize_pdf_text("\n".join(ocr_pages))
        if not merged:
            return "", "empty-text"
        return merged, ""

    def _baidu_ocr_configured(self) -> bool:
        return bool(os.getenv("BAIDU_OCR_API_KEY", "").strip() and os.getenv("BAIDU_OCR_SECRET_KEY", "").strip())

    def _baidu_refresh_token(self) -> Tuple[str, str]:
        if not self._baidu_ocr_configured():
            return "", "missing-credentials"
        now = time.time()
        if self._baidu_token and now < self._baidu_token_deadline:
            return self._baidu_token, ""
        api_key = os.getenv("BAIDU_OCR_API_KEY", "").strip()
        secret = os.getenv("BAIDU_OCR_SECRET_KEY", "").strip()
        try:
            with httpx.Client(timeout=30.0) as client:
                r = client.get(
                    "https://aip.baidubce.com/oauth/2.0/token",
                    params={
                        "grant_type": "client_credentials",
                        "client_id": api_key,
                        "client_secret": secret,
                    },
                )
        except Exception as exc:
            return "", f"token-request:{exc}"
        if r.status_code != 200:
            return "", f"token-http-{r.status_code}"
        try:
            data = r.json()
        except Exception:
            return "", "token-bad-json"
        token = data.get("access_token")
        if not token:
            err = data.get("error_description") or data.get("error") or "no-token"
            return "", f"token-error:{err}"
        self._baidu_token = str(token)
        exp = int(data.get("expires_in") or 2592000)
        self._baidu_token_deadline = now + max(120.0, float(exp) - 120.0)
        return self._baidu_token, ""

    def _image_to_baidu_jpeg_base64(self, img: Any) -> Tuple[str, str]:
        from PIL import Image as PILImage

        try:
            im = img.convert("RGB")
        except Exception as exc:
            return "", f"convert-rgb:{exc}"
        max_side = self._safe_int(os.getenv("PDF_OCR_BAIDU_MAX_SIDE", "4096"), 4096, min_value=800, max_value=4096)
        w, h = im.size
        if max(w, h) > max_side:
            scale = max_side / float(max(w, h))
            try:
                resample = PILImage.Resampling.LANCZOS
            except AttributeError:
                resample = PILImage.LANCZOS  # type: ignore[attr-defined]
            im = im.resize((max(1, int(w * scale)), max(1, int(h * scale))), resample)
        max_raw = self._safe_int(os.getenv("PDF_OCR_BAIDU_MAX_JPEG_BYTES", "3145728"), 3145728, min_value=200000, max_value=4194304)
        quality = self._safe_int(os.getenv("PDF_OCR_BAIDU_JPEG_QUALITY", "85"), 85, min_value=40, max_value=95)
        raw: bytes = b""
        for _ in range(6):
            buf = io.BytesIO()
            try:
                im.save(buf, format="JPEG", quality=quality)
            except Exception as exc:
                return "", f"jpeg-encode:{exc}"
            raw = buf.getvalue()
            if len(raw) <= max_raw or quality <= 45:
                break
            quality = max(45, quality - 12)
        b64 = base64.b64encode(raw).decode("ascii")
        if len(b64) > 4_000_000:
            return "", "image-too-large-after-base64"
        return b64, ""

    def _ocr_pdf_pages_baidu(self, file_path: str, dpi: int, max_pages: int) -> Tuple[str, str, int]:
        """返回 (text, error, successful_baidu_api_calls)。"""
        if not self._baidu_ocr_configured():
            return "", "not-configured", 0
        token, terr = self._baidu_refresh_token()
        if not token:
            return "", terr or "no-token", 0
        product = (os.getenv("BAIDU_OCR_PRODUCT", "accurate_basic").strip().lower() or "accurate_basic")
        if product not in {"accurate_basic", "general_basic"}:
            product = "accurate_basic"
        ocr_url = f"https://aip.baidubce.com/rest/2.0/ocr/v1/{product}"
        timeout = float(self._safe_int(os.getenv("PDF_OCR_BAIDU_TIMEOUT_SEC", "60"), 60, min_value=15, max_value=300))
        delay_ms = self._safe_int(os.getenv("PDF_OCR_BAIDU_PAGE_DELAY_MS", "0"), 0, min_value=0, max_value=5000)
        ocr_pages: List[str] = []
        page_idx = 0
        api_ok_count = 0
        try:
            for img in self._iter_pdf_page_images(file_path, dpi, max_pages):
                page_idx += 1
                if delay_ms > 0 and page_idx > 1:
                    time.sleep(delay_ms / 1000.0)
                try:
                    b64, enc_err = self._image_to_baidu_jpeg_base64(img)
                    if not b64:
                        return "", enc_err or "encode-failed", api_ok_count
                    with httpx.Client(timeout=timeout) as client:
                        r = client.post(
                            ocr_url,
                            params={"access_token": token},
                            data={"image": b64},
                            headers={"Content-Type": "application/x-www-form-urlencoded"},
                        )
                finally:
                    try:
                        img.close()
                    except Exception:
                        pass
                if r.status_code != 200:
                    return "", f"http-{r.status_code}", api_ok_count
                try:
                    payload = r.json()
                except Exception:
                    return "", "bad-json", api_ok_count
                err_code = payload.get("error_code")
                if err_code not in (None, 0):
                    return "", f"api-{err_code}:{payload.get('error_msg', '')}"[:200], api_ok_count
                api_ok_count += 1
                words = payload.get("words_result") or []
                lines = [str(item.get("words", "")).strip() for item in words if isinstance(item, dict)]
                block = "\n".join(lines)
                if block:
                    ocr_pages.append(block)
        except Exception as exc:
            return "", f"ocr-failed:{exc}", api_ok_count
        merged = self._normalize_pdf_text("\n".join(ocr_pages))
        if not merged:
            return "", "empty-text", api_ok_count
        return merged, "", api_ok_count

    def _get_paddle_runtime(self) -> Optional[Any]:
        if self._paddle_ocr is not None:
            return self._paddle_ocr
        if self._paddle_init_error is not None:
            return None
        if PaddleOCR is None:
            self._paddle_init_error = "paddleocr-not-installed"
            return None
        paddle_lang = (os.getenv("PADDLEOCR_LANG", "ch").strip() or "ch")
        paddle_use_angle_cls = self._as_bool(os.getenv("PADDLEOCR_USE_ANGLE_CLS", "1"), True)
        paddle_show_log = self._as_bool(os.getenv("PADDLEOCR_SHOW_LOG", "0"), False)
        # Default off: PaddlePaddle 3.3+ CPU oneDNN/PIR can raise Unimplemented in onednn_instruction.cc.
        # Optional: set PADDLEOCR_ENABLE_MKLDNN=1 to re-enable MKLDNN when your stack supports it.
        paddle_enable_mkldnn = self._as_bool(os.getenv("PADDLEOCR_ENABLE_MKLDNN", "0"), False)
        paddle_cpu_threads = self._safe_int(os.getenv("PADDLEOCR_CPU_THREADS", "1"), 1, min_value=1, max_value=16)
        # GPU 自动检测：PADDLEOCR_USE_GPU=auto（默认）时自动检测 CUDA，有 GPU 则开启；
        # 设为 1 强制开启，设为 0 强制关闭。
        gpu_env = (os.getenv("PADDLEOCR_USE_GPU", "auto") or "auto").strip().lower()
        if gpu_env == "auto":
            paddle_use_gpu = self._detect_gpu()
        else:
            paddle_use_gpu = self._as_bool(gpu_env, False)
        base_core: Dict[str, Any] = {"use_angle_cls": paddle_use_angle_cls, "lang": paddle_lang}
        base_safe: Dict[str, Any] = {**base_core, "use_gpu": paddle_use_gpu, "cpu_threads": paddle_cpu_threads}
        base_no_threads: Dict[str, Any] = {**base_core, "use_gpu": paddle_use_gpu}
        cpu_safe: Dict[str, Any] = {**base_core, "use_gpu": False, "cpu_threads": paddle_cpu_threads}
        cpu_no_threads: Dict[str, Any] = {**base_core, "use_gpu": False}
        # 先尝试不含 show_log（部分版本报 Unknown argument: show_log，且抛的是 Exception 而非 TypeError）
        attempts: List[Dict[str, Any]] = [
            {**base_safe, "enable_mkldnn": paddle_enable_mkldnn},
            {**base_safe, "show_log": paddle_show_log, "enable_mkldnn": paddle_enable_mkldnn},
            {**base_safe, "show_log": paddle_show_log},
            dict(base_safe),
            {**base_no_threads, "enable_mkldnn": paddle_enable_mkldnn},
            {**base_no_threads, "show_log": paddle_show_log, "enable_mkldnn": paddle_enable_mkldnn},
            {**base_no_threads, "show_log": paddle_show_log},
            dict(base_no_threads),
            {**cpu_safe, "enable_mkldnn": paddle_enable_mkldnn},
            {**cpu_safe, "show_log": paddle_show_log, "enable_mkldnn": paddle_enable_mkldnn},
            {**cpu_safe, "show_log": paddle_show_log},
            dict(cpu_safe),
            {**cpu_no_threads, "enable_mkldnn": paddle_enable_mkldnn},
            {**cpu_no_threads, "show_log": paddle_show_log, "enable_mkldnn": paddle_enable_mkldnn},
            {**cpu_no_threads, "show_log": paddle_show_log},
            dict(cpu_no_threads),
            {**base_core, "enable_mkldnn": paddle_enable_mkldnn},
            {**base_core, "show_log": paddle_show_log, "enable_mkldnn": paddle_enable_mkldnn},
            {**base_core, "show_log": paddle_show_log},
            dict(base_core),
        ]
        last_error: Optional[Exception] = None
        for kwargs in attempts:
            try:
                self._paddle_ocr = PaddleOCR(**kwargs)
                return self._paddle_ocr
            except TypeError as exc:
                last_error = exc
                continue
            except Exception as exc:
                msg = str(exc).lower()
                if (
                    "show_log" in msg
                    or "unknown argument" in msg
                    or "unexpected keyword" in msg
                    or "got an unexpected keyword" in msg
                ):
                    last_error = exc
                    continue
                self._paddle_init_error = str(exc)
                return None
        self._paddle_init_error = str(last_error) if last_error else "PaddleOCR init failed"
        return None

    def _extract_paddle_texts(self, result: Any) -> List[str]:
        texts: List[str] = []
        if not isinstance(result, list):
            return texts
        for page in result:
            if not isinstance(page, list):
                continue
            for line in page:
                if isinstance(line, (list, tuple)) and len(line) >= 2:
                    content = line[1]
                    if isinstance(content, (list, tuple)) and content:
                        text = str(content[0]).strip()
                        if text:
                            texts.append(text)
        return texts

    def _parse_pptx(self, file_path: str) -> str:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Pt  # type: ignore  # noqa: F401
        except ImportError as exc:
            raise ValueError("需要安装 python-pptx：pip install python-pptx") from exc
        try:
            prs = Presentation(file_path)
            slides: List[str] = []
            for i, slide in enumerate(prs.slides, 1):
                lines: List[str] = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            lines.append(t)
                if lines:
                    slides.append(f"[幻灯片 {i}]\n" + "\n".join(lines))
            return "\n\n".join(slides).strip()
        except Exception as exc:
            raise ValueError(f"PPTX解析失败: {exc}") from exc

    def _parse_image_file(self, file_path: str) -> str:
        """对独立图片文件（PNG/JPG 等）直接做 OCR，复用三引擎级联逻辑。"""
        from PIL import Image as PILImage  # type: ignore

        try:
            img = PILImage.open(file_path)
        except Exception as exc:
            raise ValueError(f"图片文件打开失败: {exc}") from exc

        ocr_dpi = self._safe_int(os.getenv("PDF_OCR_DPI", "260"), 260, min_value=120, max_value=500)
        errors: List[str] = []

        # 尝试 PaddleOCR
        if PaddleOCR is not None and np is not None:
            runtime = self._get_paddle_runtime()
            if runtime is not None:
                try:
                    arr = np.array(img.convert("RGB"))
                    try:
                        result = runtime.ocr(arr, cls=True)
                    except TypeError:
                        result = runtime.ocr(arr)
                    texts = self._extract_paddle_texts(result)
                    if texts:
                        img.close()
                        return "\n".join(texts)
                except Exception as exc:
                    errors.append(f"paddle:{exc}")

        # 尝试 Tesseract
        if pytesseract is not None:
            try:
                ocr_lang = os.getenv("PDF_OCR_TESSERACT_LANG", "chi_sim+eng").strip() or "chi_sim+eng"
                text = pytesseract.image_to_string(img, lang=ocr_lang)
                if text and text.strip():
                    img.close()
                    return text.strip()
            except Exception as exc:
                errors.append(f"tesseract:{exc}")

        # 尝试百度 OCR
        if self._baidu_ocr_configured():
            try:
                b64, enc_err = self._image_to_baidu_jpeg_base64(img)
                if b64:
                    token, token_err = self._baidu_refresh_token()
                    if token:
                        ocr_url = os.getenv("BAIDU_OCR_URL", "https://aip.baidubce.com/rest/2.0/ocr/v1/accurate_basic").strip()
                        with httpx.Client(timeout=60.0) as client:
                            r = client.post(
                                ocr_url,
                                params={"access_token": token},
                                data={"image": b64},
                                headers={"Content-Type": "application/x-www-form-urlencoded"},
                            )
                        data = r.json()
                        words = [w.get("words", "") for w in data.get("words_result", [])]
                        text = "\n".join(words)
                        if text.strip():
                            img.close()
                            return text.strip()
            except Exception as exc:
                errors.append(f"baidu:{exc}")

        img.close()
        if errors:
            raise ValueError(f"图片OCR失败: {'; '.join(errors)}")
        raise ValueError("图片OCR失败：未安装任何OCR引擎（paddleocr/pytesseract）或百度OCR未配置")

    def _parse_docx(self, file_path: str) -> str:
        try:
            doc = Document(file_path)
            parts: List[str] = []

            # 正文段落
            for p in doc.paragraphs:
                t = p.text.strip()
                if t:
                    parts.append(t)

            # 表格（逐行逐格提取，用制表符分隔列、换行分隔行）
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append("\t".join(cells))

            # 页眉页脚（各节）
            for section in doc.sections:
                for hf in (section.header, section.footer):
                    try:
                        for p in hf.paragraphs:
                            t = p.text.strip()
                            if t:
                                parts.append(t)
                    except Exception:
                        pass

            return "\n".join(parts).strip()
        except Exception as exc:
            raise ValueError(f"DOCX解析失败: {exc}") from exc

    def _parse_text(self, file_path: str) -> str:
        for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk", "latin-1"):
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except Exception:
                continue
        raise ValueError("文本文件编码无法识别")

    def _normalize_pdf_text(self, text: str) -> str:
        cleaned = (text or "").replace("\x00", " ").replace("\ufeff", "")
        cleaned = re.sub(r"[ \t]+", " ", cleaned)
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
        return cleaned.strip()

    def _prefer_direct_pdf_text(self, text: str) -> bool:
        normalized = (text or "").strip()
        if not normalized:
            return False
        signal_chars = len(re.findall(r"[\u4e00-\u9fffA-Za-z0-9]", normalized))
        score = self._text_quality_score(normalized)
        if signal_chars >= 120:
            return True
        return signal_chars >= 30 and score >= 0.45

    def _text_quality_score(self, text: str) -> float:
        if not text:
            return 0.0
        total = len(text)
        cjk = len(re.findall(r"[\u4e00-\u9fff]", text))
        alpha = len(re.findall(r"[A-Za-z]", text))
        replacement = text.count("\ufffd") + text.count("?")
        control_like = len(re.findall(r"[\x01-\x08\x0B\x0C\x0E-\x1F]", text))
        # Prefer content with more meaningful characters and fewer malformed symbols.
        signal = cjk + alpha
        noise = replacement + control_like
        return (signal / max(total, 1)) - (noise / max(total, 1))

    def _safe_int(self, value: str, default: int, min_value: int, max_value: int) -> int:
        try:
            parsed = int(str(value).strip())
            return max(min_value, min(max_value, parsed))
        except Exception:
            return default

    def _detect_gpu(self) -> bool:
        """自动检测是否有可用 GPU（CUDA 或 ROCm），有则返回 True。"""
        try:
            import paddle  # type: ignore
            dev = paddle.device.get_device()
            return str(dev).startswith("gpu")
        except Exception:
            pass
        try:
            import torch  # type: ignore
            return torch.cuda.is_available()
        except Exception:
            pass
        return False

    def _as_bool(self, value: str, default: bool) -> bool:
        normalized = str(value).strip().lower()
        if normalized in {"1", "true", "yes", "on"}:
            return True
        if normalized in {"0", "false", "no", "off"}:
            return False
        return default

    def _extract_metadata(self, text: str, document_type: str, filename: str) -> Dict[str, Any]:
        first_lines = [line.strip() for line in text.split("\n") if line.strip()][:8]
        title = first_lines[0] if first_lines else filename

        headings = self._extract_headings(text)
        knowledge_points = self._extract_knowledge_points(text, document_type)
        discipline = self._infer_discipline(text, filename)

        return {
            "title": title[:180],
            "filename": filename,
            "document_type": document_type,
            "discipline": discipline,
            "headings": headings[:30],
            "knowledge_points": knowledge_points[:60],
            "length": len(text),
        }

    def _extract_headings(self, text: str) -> List[str]:
        candidates = []
        for line in text.split("\n"):
            clean = line.strip()
            if not clean:
                continue
            if re.match(r"^(\d+(\.\d+)*)\s+.+", clean) or clean.endswith(":") or clean.endswith("："):
                candidates.append(clean)
            elif len(clean) <= 35 and clean.isupper():
                candidates.append(clean)
        return candidates

    def _extract_knowledge_points(self, text: str, document_type: str) -> List[str]:
        kws = set()
        if document_type == "exam":
            patterns = [r"考点[:：]\s*([^\n]+)", r"知识点[:：]\s*([^\n]+)"]
        elif document_type == "technical":
            patterns = [r"API[:：]\s*([A-Za-z0-9_/\-]+)", r"endpoint[:：]\s*([A-Za-z0-9_/\-]+)"]
        elif document_type == "project":
            patterns = [r"任务[:：]\s*([^\n]+)", r"里程碑[:：]\s*([^\n]+)"]
        else:
            patterns = [r"关键词[:：]\s*([^\n]+)", r"关键字[:：]\s*([^\n]+)"]

        for p in patterns:
            for match in re.findall(p, text, flags=re.IGNORECASE):
                for item in re.split(r"[，,;；、]\s*", match):
                    i = item.strip()
                    if 1 < len(i) <= 40:
                        kws.add(i)

        if not kws:
            words = re.findall(r"[A-Za-z]{4,}|[\u4e00-\u9fff]{2,8}", text)
            for w in words[:300]:
                if len(w) >= 2:
                    kws.add(w.lower() if re.match(r"[A-Za-z]+$", w) else w)
        return list(kws)

    def _infer_discipline(self, text: str, filename: str) -> str:
        corpus = f"{filename} {text[:2000]}".lower()
        mapping = {
            "computer-science": ["algorithm", "api", "database", "模型", "算法", "编程", "系统"],
            "mathematics": ["theorem", "proof", "matrix", "概率", "统计", "线性代数"],
            "physics": ["quantum", "mechanics", "entropy", "力学", "电磁", "热力学"],
            "economics": ["market", "finance", "gdp", "经济", "金融", "博弈"],
            "medicine": ["clinical", "patient", "treatment", "医学", "病理", "临床"],
            "general": [],
        }
        best_name = "general"
        best_score = 0
        for name, terms in mapping.items():
            score = sum(1 for t in terms if t in corpus)
            if score > best_score:
                best_score = score
                best_name = name
        return best_name

    def _infer_document_form(self, text: str, filename: str, page_count: int = 0) -> str:
        """
        自动识别文档形态，返回：
        paper（学术论文）/ thesis（学位论文）/ textbook（教材）/ book（书籍）/
        exam（试卷）/ report（报告）/ slides（幻灯片）/ project（项目文档）/
        technical（技术文档）/ general（通用）
        """
        corpus = (filename + " " + text[:4000]).lower()
        tail = text[-2000:].lower() if len(text) > 2000 else corpus
        scores: dict = {k: 0.0 for k in (
            "paper", "thesis", "textbook", "book", "exam",
            "report", "slides", "project", "technical", "general"
        )}
        # 学位论文
        for kw in ["学位论文", "毕业论文", "硕士论文", "博士论文",
                   "dissertation", "master thesis", "phd thesis", "致谢", "acknowledgements"]:
            if kw in corpus:
                scores["thesis"] += 2
        # 学术论文
        for kw in ["abstract", "摘要", "keywords", "关键词", "doi:", "issn",
                   "journal", "期刊", "vol.", "volume", "references", "参考文献", "bibliography"]:
            if kw in corpus or kw in tail:
                scores["paper"] += 1.5
        if re.search(r"\[\d+\]|\(\w+,\s*\d{4}\)", text):
            scores["paper"] += 2
        # 教材
        for kw in [r"第.{1,3}章", "习题", "练习题", "本章小结", "学习目标", "课后", "exercise"]:
            if re.search(kw, corpus):
                scores["textbook"] += 1.5
        # 书籍
        for kw in ["isbn", "出版社", "版权所有", "publisher", "copyright", "前言", "preface", "目录", "contents"]:
            if kw in corpus:
                scores["book"] += 1.2
        if page_count > 80:
            scores["book"] += 1
        if page_count > 200:
            scores["book"] += 1
        # 试卷
        for kw in ["试卷", "考试", "满分", "答题", "考生", "exam paper", "time allowed", "total marks"]:
            if kw in corpus:
                scores["exam"] += 2
        if re.search(r"(第\s*[一二三四五六七八九十\d]+\s*题|^\s*\d+[\.、。])", text, re.MULTILINE):
            scores["exam"] += 1
        # 报告
        for kw in ["调研报告", "研究报告", "分析报告", "工作报告",
                   "executive summary", "findings", "recommendations", "结论与建议"]:
            if kw in corpus:
                scores["report"] += 2
        # 幻灯片
        lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
        short_ratio = sum(1 for ln in lines if len(ln) < 30) / max(len(lines), 1)
        if short_ratio > 0.65 and len(lines) > 10:
            scores["slides"] += 2
        if re.search(r"slide|ppt|pptx", corpus):
            scores["slides"] += 2
        # 项目文档
        for kw in ["需求文档", "设计文档", "需求分析", "架构设计", "里程碑", "甘特图",
                   "user story", "sprint", "backlog"]:
            if kw in corpus:
                scores["project"] += 2
        # 技术文档
        for kw in ["readme", "installation", "api reference", "endpoint",
                   "configuration", "getting started", "安装", "配置说明"]:
            if kw in corpus:
                scores["technical"] += 1.5
        if re.search(r"```|`[^`]+`", text):
            scores["technical"] += 1
        # thesis 优先级高于 paper
        if scores["thesis"] > 0:
            scores["paper"] *= 0.5
        best = max(scores, key=lambda k: scores[k])
        return best if scores[best] > 0 else "general"

    def _detect_has_images(self, file_path: str, ext: str) -> "tuple[bool, int]":
        """检测文档中是否包含图片，返回 (has_images, image_count_estimate)。"""
        if ext == ".pdf":
            return self._detect_pdf_images(file_path)
        if ext == ".docx":
            return self._detect_docx_images(file_path)
        return False, 0

    def _detect_pdf_images(self, file_path: str) -> "tuple[bool, int]":
        count = 0
        try:
            reader = PdfReader(file_path)
            for page in reader.pages:
                resources = page.get("/Resources")
                if not resources:
                    continue
                xobjects = resources.get("/XObject")
                if not xobjects:
                    continue
                try:
                    xobj_dict = xobjects.get_object() if hasattr(xobjects, "get_object") else xobjects
                    for _, obj_ref in xobj_dict.items():
                        try:
                            obj = obj_ref.get_object() if hasattr(obj_ref, "get_object") else obj_ref
                            if obj.get("/Subtype") == "/Image":
                                count += 1
                        except Exception:
                            pass
                except Exception:
                    pass
        except Exception:
            pass
        if count == 0 and pdfplumber is not None:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        count += len(page.images or [])
            except Exception:
                pass
        return count > 0, count

    def _detect_docx_images(self, file_path: str) -> "tuple[bool, int]":
        count = 0
        try:
            doc = Document(file_path)
            count = len(doc.inline_shapes)
            if count == 0:
                for rel in doc.part.rels.values():
                    if "image" in str(rel.reltype).lower():
                        count += 1
        except Exception:
            pass
        return count > 0, count

    def _check_encoding_quality(self, text: str) -> dict:
        """检测文本编码质量，返回 {"ok": bool, "issues": List[str], "garble_ratio": float}。"""
        issues: list = []
        if not text:
            return {"ok": True, "issues": [], "garble_ratio": 0.0}
        total = len(text)
        # Unicode 替换字符
        replacement_ratio = text.count("�") / total
        if replacement_ratio > 0.02:
            issues.append(f"replacement_char_ratio={replacement_ratio:.2%}")
        # 不可打印控制字符（换行/制表符除外）
        ctrl_count = sum(1 for c in text if ord(c) < 32 and c not in "\n\r\t")
        ctrl_ratio = ctrl_count / total
        if ctrl_ratio > 0.01:
            issues.append(f"control_char_ratio={ctrl_ratio:.2%}")
        # Latin-1 被误读为 UTF-8 的典型乱码模式
        mojibake_hits = len(re.findall(r"[Ã-ÿ]{2,}|â€[’“”]", text))
        if mojibake_hits > 5:
            issues.append(f"mojibake_patterns={mojibake_hits}")
        # 大量非 CJK 高位字符
        non_cjk_high = sum(
            1 for c in text
            if ord(c) > 127
            and not ("一" <= c <= "鿿")
            and not ("　" <= c <= "〿")
            and not ("＀" <= c <= "￯")
        )
        if total > 200 and non_cjk_high / total > 0.15:
            issues.append(f"high_non_cjk_ratio={non_cjk_high / total:.2%}")
        garble_ratio = max(replacement_ratio, ctrl_ratio)
        return {"ok": len(issues) == 0, "issues": issues, "garble_ratio": round(garble_ratio, 4)}
